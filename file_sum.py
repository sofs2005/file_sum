# encoding:utf-8
import json
import os
import csv
import re
import requests
import plugins
from bridge.context import ContextType
from bridge.reply import Reply, ReplyType
from common.log import logger
from common.expired_dict import ExpiredDict
from plugins import *
from docx import Document
from bs4 import BeautifulSoup
from pptx import Presentation
from openpyxl import load_workbook
import fitz  # PyMuPDF
import subprocess

# 文件扩展名到类型的映射
EXTENSION_TO_TYPE = {
    'pdf': 'pdf',
    'docx': 'docx',
    'doc': 'docx',
    'md': 'md',
    'markdown': 'md',
    'txt': 'txt',
    'xlsx': 'excel',
    'xls': 'excel',
    'csv': 'csv',
    'html': 'html',
    'htm': 'html',
    'pptx': 'ppt',
    'ppt': 'ppt'
}

@plugins.register(
    name="FileSum",
    desire_priority=20,
    hidden=False,
    desc="A plugin for summarizing files",
    version="1.1.0",
    author="sofs2005",
)
class FileSum(Plugin):
    def __init__(self):
        super().__init__()
        try:
            # 加载配置
            self.config = super().load_config()
            if not self.config:
                self.config = self._load_config_template()
            
            # 初始化配置
            self.max_file_size = self.config.get("max_file_size", 15000)
            self.max_token_size = self.config.get("max_token_size", 4000)
            self.group = self.config.get("group", True)
            self.prompt = self.config.get("prompt", "请总结这个文件的主要内容")
            
            # 初始化缓存
            self.file_cache = ExpiredDict(self.config.get("file_cache_time", 300))
            
            # 注册事件处理器
            self.handlers[Event.ON_HANDLE_CONTEXT] = self.on_handle_context
            
            logger.info("[FileSum] Plugin initialized")
        except Exception as e:
            logger.error(f"[FileSum] Init failed: {e}")
            raise e

    def get_help_text(self, **kwargs):
        help_text = "📄 文件总结插件使用说明：\n"
        help_text += "1. 发送文件后，单聊会自动总结\n"
        help_text += "2. 群聊需要发送「总结」触发总结\n"
        help_text += "\n支持格式：PDF、Word、Excel、PPT、TXT、Markdown、HTML、CSV"
        return help_text

    def _load_config_template(self):
        try:
            plugin_config_path = os.path.join(self.path, "config.json.template")
            if os.path.exists(plugin_config_path):
                with open(plugin_config_path, "r", encoding="utf-8") as f:
                    return json.load(f)
        except Exception as e:
            logger.error(f"[FileSum] Load config template failed: {e}")
        return {}

    def on_handle_context(self, e_context: EventContext):
        context = e_context["context"]
        msg: ChatMessage = e_context["context"]["msg"]
        
        # 生成缓存键
        chat_id = context.get("session_id", "default")
        user_id = msg.from_user_id
        
        # 清理ID中的特殊字符
        chat_id = chat_id.replace('@', '').split('_')[0]
        user_id = user_id.replace('@', '').split('_')[0]
        
        # 生成缓存key
        cache_key = f"filesum_{chat_id}_{user_id}"

        isgroup = e_context["context"].get("isgroup", False)
        
        if isgroup and not self.group:
            logger.info("[FileSum] 群聊消息，文件处理功能已禁用")
            return

        # 处理文件消息
        if context.type == ContextType.FILE:
            logger.info(f"[FileSum] 收到文件，存入缓存，key={cache_key}")
            context.get("msg").prepare()
            file_path = context.content
            
            self.file_cache[cache_key] = {
                'file_path': file_path,
                'processed': False
            }

            # 如果是单聊，直接触发总结
            if not isgroup:
                logger.info("[FileSum] 单聊消息，自动触发总结")
                return self._process_file_summary(cache_key, e_context)
            return

        # 处理文本消息
        if context.type == ContextType.TEXT:
            text = context.content
            
            # 群聊中的总结触发命令
            if isgroup and "总结" in text:
                logger.info("[FileSum] 群聊中收到总结命令")
                if cache_key in self.file_cache:
                    logger.info(f"[FileSum] 找到文件缓存，开始处理总结")
                    return self._process_file_summary(cache_key, e_context)
                else:
                    logger.info("[FileSum] 未找到待处理的文件，让事件继续传递")
                    return False  # 返回 False 让事件继续传递给其他插件（如 JinaSum）

        return False

    def _process_file_summary(self, cache_key: str, e_context: EventContext):
        """处理文件总结的核心逻辑"""
        try:
            cache_data = self.file_cache.get(cache_key)
            if not cache_data:
                logger.info("[filesum] 未找到缓存的文件")
                return
            
            file_path = cache_data.get('file_path')
            if not file_path or not os.path.exists(file_path):
                logger.info("[filesum] 缓存的文件不存在")
                reply = Reply(ReplyType.ERROR, "文件已过期，请重新发送")
                e_context["reply"] = reply
                e_context.action = EventAction.BREAK_PASS
                return

            # 发送处理中的提示
            reply = Reply(ReplyType.TEXT, "📄 正在为您总结文件内容，请稍候...")
            channel = e_context["channel"]
            channel.send(reply, e_context["context"])

            # 读取文件内容
            logger.info("[filesum] 开始读取文件内容")
            file_content = self.extract_content(file_path)
            if file_content is None:
                logger.info("[filesum] 文件内容无法提取")
                reply = Reply(ReplyType.ERROR, "无法读取文件内容")
                e_context["reply"] = reply
                e_context.action = EventAction.BREAK_PASS
                return

            # 处理文件内容
            self.handle_file(file_content, e_context)
            
            # 处理完成后删除文件
            try:
                os.remove(file_path)
                logger.info(f"[filesum] 文件 {file_path} 已删除")
                # 删除文件路径缓存
                del self.file_cache[cache_key]
                # 设置事件状态为 CONTINUE
                e_context.action = EventAction.CONTINUE
                return True
            except Exception as e:
                logger.error(f"[filesum] 删除文件失败: {str(e)}")
                return False

        except Exception as e:
            logger.error(f"[filesum] 处理文件总结时出错: {str(e)}")
            reply = Reply(ReplyType.ERROR, f"处理文件时出错: {str(e)}")
            e_context["reply"] = reply
            e_context.action = EventAction.BREAK_PASS
            return False

    def extract_content(self, file_path):
        """提取文件内容"""
        try:
            # 添加文件大小检查
            file_size = os.path.getsize(file_path) / 1024  # 转换为KB
            if file_size > self.max_file_size:
                logger.error(f"文件大小 ({file_size}KB) 超过限制 ({self.max_file_size}KB)")
                return None
            
            file_extension = os.path.splitext(file_path)[1].lower().replace('.', '')
            file_type = EXTENSION_TO_TYPE.get(file_extension)
            
            if file_type == 'pdf':
                return self.read_pdf(file_path)
            elif file_type == 'docx':
                return self.read_docx(file_path)
            elif file_type == 'md':
                return self.read_markdown(file_path)
            elif file_type == 'txt':
                return self.read_txt(file_path)
            elif file_type == 'excel':
                return self.read_excel(file_path)
            elif file_type == 'csv':
                return self.read_csv(file_path)
            elif file_type == 'html':
                return self.read_html(file_path)
            elif file_type == 'ppt':
                return self.read_ppt(file_path)
            else:
                logger.error(f"不支持的文件类型: {file_extension}")
                return None
        except Exception as e:
            logger.error(f"提取文件内容时出错: {str(e)}")
            return None

    def read_pdf(self, file_path):
        """读取PDF文件"""
        try:
            doc = fitz.open(file_path)
            content = ' '.join([page.get_text() for page in doc])
            doc.close()
            return content
        except Exception as e:
            logger.error(f"读取PDF文件失败: {str(e)}")
            return None

    def read_docx(self, file_path):
        """读取Word文档"""
        try:
            # 检查文件扩展名
            ext = os.path.splitext(file_path)[1].lower()
            
            if ext == '.docx':
                # 处理 .docx 文件
                doc = Document(file_path)
                content = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
                return content
            elif ext == '.doc':
                # 处理 .doc 文件，使用系统 antiword 命令
                try:
                    # 使用绝对路径
                    antiword_path = '/usr/bin/antiword'
                    if not os.path.exists(antiword_path):
                        logger.error("未找到 antiword，请先安装: sudo apt-get install antiword")
                        return None
                        
                    result = subprocess.run([antiword_path, file_path], 
                                         capture_output=True, 
                                         text=True,
                                         encoding='utf-8')
                    if result.returncode == 0:
                        return result.stdout
                    else:
                        logger.error(f"antiword 处理失败: {result.stderr}")
                        return None
                except Exception as e:
                    logger.error(f"使用 antiword 处理文件失败: {str(e)}")
                    return None
            else:
                logger.error(f"不支持的Word文件格式: {ext}")
                return None
            
        except Exception as e:
            logger.error(f"读取Word文档失败: {str(e)}")
            return None

    def read_markdown(self, file_path):
        """读取Markdown文件"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            return remove_markdown(content)
        except Exception as e:
            logger.error(f"读取Markdown文件失败: {str(e)}")
            return None

    def read_txt(self, file_path):
        """读取文本文件"""
        encodings = ['utf-8', 'gbk', 'gb2312', 'ascii']
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
            except Exception as e:
                logger.error(f"读取文本文件失败: {str(e)}")
                return None
        return None

    def read_excel(self, file_path):
        """读取Excel文件"""
        try:
            ext = os.path.splitext(file_path)[1].lower()
            
            if ext == '.xlsx':
                # 处理 .xlsx 文件
                wb = load_workbook(file_path)
                content = []
                for sheet in wb.worksheets:
                    for row in sheet.iter_rows(values_only=True):
                        content.append('\t'.join([str(cell) if cell is not None else '' for cell in row]))
                return '\n'.join(content)
            elif ext == '.xls':
                # 处理 .xls 文件
                try:
                    import xlrd
                    wb = xlrd.open_workbook(file_path)
                    content = []
                    for sheet in wb.sheets():
                        for row in range(sheet.nrows):
                            row_values = [str(cell.value) if cell.value is not None else '' for cell in sheet.row(row)]
                            content.append('\t'.join(row_values))
                    return '\n'.join(content)
                except ImportError:
                    logger.error("未安装 xlrd 库，无法读取 .xls 文件。请安装：pip install xlrd")
                    return None
            else:
                logger.error(f"不支持的Excel文件格式: {ext}")
                return None
        except Exception as e:
            logger.error(f"读取Excel文件失败: {str(e)}")
            return None

    def read_csv(self, file_path):
        """读取CSV文件"""
        try:
            content = []
            with open(file_path, 'r', encoding='utf-8') as f:
                reader = csv.reader(f)
                for row in reader:
                    content.append('\t'.join(row))
            return '\n'.join(content)
        except Exception as e:
            logger.error(f"读取CSV文件失败: {str(e)}")
            return None

    def read_html(self, file_path):
        """读取HTML文件"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                soup = BeautifulSoup(f.read(), 'html.parser')
                return soup.get_text()
        except Exception as e:
            logger.error(f"读取HTML文件失败: {str(e)}")
            return None

    def read_ppt(self, file_path):
        """读取PPT文件"""
        try:
            ext = os.path.splitext(file_path)[1].lower()
            
            if ext == '.pptx':
                # 处理 .pptx 文件
                prs = Presentation(file_path)
                content = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            content.append(shape.text)
                return '\n'.join(content)
            elif ext == '.ppt':
                logger.error("不支持旧版 .ppt 格式，请转换为 .pptx 后重试")
                return None
            else:
                logger.error(f"不支持的PPT文件格式: {ext}")
                return None
        except Exception as e:
            logger.error(f"读取PPT文件失败: {str(e)}")
            return None

    def handle_file(self, content, e_context):
        """处理文件内容"""
        try:
            if not content:
                reply = Reply(ReplyType.ERROR, "无法读取文件内容")
                e_context["reply"] = reply
                return

            # 用置中的token限制
            if len(content) > self.max_token_size:
                content = content[:self.max_token_size] + "..."
                logger.warning(f"文件内容已截断到 {self.max_token_size} 个字符")

            # 构建提示词，移除追问相关的提示
            prompt = f"{self.prompt}\n\n{content}"
            
            # 设置用户消息
            e_context["context"].type = ContextType.TEXT
            e_context["context"].content = prompt
            
            # 让事件继续传递给 bot 处理
            e_context.action = EventAction.CONTINUE
            return True

        except Exception as e:
            logger.error(f"处理文件内容时出错: {str(e)}")
            reply = Reply(ReplyType.ERROR, f"处理文件时出错: {str(e)}")
            e_context["reply"] = reply
            e_context.action = EventAction.BREAK_PASS

def remove_markdown(text):
    """移除Markdown格式"""
    # 移除标题
    text = re.sub(r'#{1,6}\s+', '', text)
    # 移除加粗和斜体
    text = re.sub(r'\*{1,2}(.*?)\*{1,2}', r'\1', text)
    # 移除链接
    text = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', text)
    # 移除代码块
    text = re.sub(r'```[\s\S]*?```', '', text)
    # 移除行内代码
    text = re.sub(r'`([^`]+)`', r'\1', text)
    return text.strip()